"""Generate a probe deck for image+content layouts.

This helper writes visible marker text into text-capable placeholders so the
visual role of each placeholder can be confirmed in PowerPoint.
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any

from pptx import Presentation


PROBE_LAYOUT_NAMES = [
    "Narrow Image and Content",
    "Wide Image and Content",
]


def _clear_existing_slides(prs: Any) -> None:
    slide_ids = list(prs.slides._sldIdLst)  # type: ignore[attr-defined]
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)  # type: ignore[attr-defined]


def build_probe(template_path: Path, output_path: Path) -> Path:
    prs = Presentation(str(template_path))
    if len(prs.slides) > 0:
        _clear_existing_slides(prs)

    layouts_by_name = {layout.name: layout for layout in prs.slide_layouts}
    for layout_name in PROBE_LAYOUT_NAMES:
        if layout_name not in layouts_by_name:
            continue

        slide = prs.slides.add_slide(layouts_by_name[layout_name])
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            ptype = str(ph.placeholder_format.type)
            name = getattr(ph, "name", "")
            marker = f"IDX {idx}\n{ptype}\n{name}"
            try:
                ph.text = marker
            except Exception:
                # Ignore non-text placeholders (e.g. picture).
                pass

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path.resolve()


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate probe deck for image-content layouts.")
    parser.add_argument("--template", required=True, help="Path to source template (.pptx).")
    parser.add_argument("--output", required=True, help="Output path for probe deck (.pptx).")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    output = build_probe(Path(args.template), Path(args.output))
    print(f"Image-content probe deck generated: {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
