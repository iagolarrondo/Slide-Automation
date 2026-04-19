"""Inspect donor deck slides and materialized placeholder/shape targets.

Usage:
    python src/inspect_template.py --template templates/Sloan_Donor_Deck.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation

from slide_automation.template_map import SLIDE_TYPE_MAP


SUPPORTED_DONOR_TYPES = [
    "cover",
    "agenda",
    "divider",
    "standard_1_block",
    "standard_2_block",
    "standard_3_block",
    "standard_2_block_big_left",
    "standard_2_block_big_right",
    "narrow_image_content",
    "wide_image_content",
    "quote",
]


def parse_args() -> argparse.Namespace:
    """Parse CLI args for donor deck inspection."""
    parser = argparse.ArgumentParser(
        description="Print donor slide inventory and shape/placeholder details."
    )
    parser.add_argument(
        "--template",
        required=True,
        help="Path to donor .pptx file.",
    )
    return parser.parse_args()


def _shape_text(shape: Any) -> str:
    if hasattr(shape, "text_frame"):
        text = (shape.text_frame.text or "").strip()
        if text:
            return text.replace("\n", " ")[:120]
    return ""


def _print_slide_inventory(prs: Any) -> None:
    print("Donor slide inventory")
    print("=" * 72)
    print(f"Total donor slides: {len(prs.slides)}")
    print("-" * 72)

    for slide_no, slide in enumerate(prs.slides, start=1):
        layout_name = str(slide.slide_layout.name)
        title_text = ""
        if getattr(slide.shapes, "title", None) is not None:
            try:
                title_text = str(slide.shapes.title.text).strip()  # type: ignore[union-attr]
            except Exception:
                title_text = ""

        mapped_types = []
        for slide_type, config in SLIDE_TYPE_MAP.items():
            if config.get("donor_slide_number") == slide_no:
                mapped_types.append(slide_type)

        candidate = ", ".join(mapped_types) if mapped_types else "unmapped"
        print(
            f"Slide {slide_no:>2}: candidate={candidate:<32} "
            f"layout={layout_name!r} title={title_text!r}"
        )

    print("-" * 72)

    missing = [
        slide_type
        for slide_type in SUPPORTED_DONOR_TYPES
        if slide_type not in SLIDE_TYPE_MAP
    ]
    if missing:
        print(f"Missing donor mapping entries: {', '.join(missing)}")
    else:
        print("All expected donor mappings are present.")
    print("=" * 72)


def _print_slide_details(prs: Any) -> None:
    for slide_no, slide in enumerate(prs.slides, start=1):
        print(f"\nSlide {slide_no}")
        print(f"  layout: {slide.slide_layout.name!r}")
        print(f"  shape_count: {len(slide.shapes)}")
        print(f"  placeholder_count: {len(slide.placeholders)}")

        print("  placeholders:")
        if not slide.placeholders:
            print("    (none materialized)")
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            ph_type = str(ph.placeholder_format.type)
            marker = _shape_text(ph)
            print(
                f"    - idx={idx:<3} type={ph_type:<18} "
                f"name={ph.name!r} marker={marker!r}"
            )

        print("  shapes:")
        for shape_idx, shape in enumerate(slide.shapes):
            marker = _shape_text(shape)
            shape_type = str(shape.shape_type)
            try:
                ph_idx = shape.placeholder_format.idx
                ph_type = str(shape.placeholder_format.type)
                ph_info = f"placeholder idx={ph_idx} type={ph_type}"
            except Exception:
                ph_info = "non-placeholder"
            print(
                f"    - shape_index={shape_idx:<2} name={shape.name!r} "
                f"type={shape_type} {ph_info} marker={marker!r}"
            )


def main() -> int:
    """Inspect donor deck and print slide-centric mapping details."""
    args = parse_args()
    template_path = Path(args.template)

    if not template_path.exists():
        print(f"Template not found: {template_path}", file=sys.stderr)
        return 1

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        print(f"Failed to open template: {exc}", file=sys.stderr)
        return 1

    print(f"Donor deck: {template_path}")
    _print_slide_inventory(prs)
    _print_slide_details(prs)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
