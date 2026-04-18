"""Generate a probe deck to visually identify placeholder roles.

This is a temporary diagnostic helper that leaves normal rendering untouched.
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any, Dict

from pptx import Presentation

from slide_automation.template_map import SLIDE_TYPE_MAP


PROBE_SLIDE_TYPES = [
    "standard_1_block",
    "standard_2_block",
    "standard_3_block",
    "standard_2_block_big_left",
    "standard_2_block_big_right",
]


def _clear_existing_slides(prs: Any) -> None:
    """Remove pre-existing slides from template-backed presentations."""
    slide_ids = list(prs.slides._sldIdLst)  # type: ignore[attr-defined]
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)  # type: ignore[attr-defined]


def _placeholder_meta(ph: Any) -> tuple[str, str]:
    """Return readable placeholder type/name metadata."""
    ptype = str(ph.placeholder_format.type)
    name = getattr(ph, "name", "")
    return ptype, name


def _is_text_role(logical_field: str) -> bool:
    """Limit probe text to title/body-like logical fields."""
    return any(
        token in logical_field
        for token in ("title", "body", "section", "footer")
    )


def _set_probe_text(slide: Any, idx: int, logical_field: str) -> None:
    """Write marker text to a mapped placeholder if writable."""
    try:
        ph = slide.placeholders[idx]
    except Exception:
        return

    ptype, name = _placeholder_meta(ph)
    marker = f"IDX {idx} | {logical_field}\n{ptype} | {name}"
    try:
        ph.text = marker
    except Exception:
        # Skip non-text placeholders.
        return


def build_probe_deck(template_path: Path, output_path: Path) -> Path:
    """Create probe deck with marker text for mapped standard placeholders."""
    prs = Presentation(str(template_path))
    if len(prs.slides) > 0:
        _clear_existing_slides(prs)

    for slide_type in PROBE_SLIDE_TYPES:
        mapping: Dict[str, Any] = SLIDE_TYPE_MAP[slide_type]
        layout_index = mapping["layout_index"]
        placeholders = mapping["placeholders"]

        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)

        for logical_field, idx in placeholders.items():
            if not isinstance(idx, int) or not _is_text_role(logical_field):
                continue
            _set_probe_text(slide, idx, logical_field)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path.resolve()


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate placeholder probe deck.")
    parser.add_argument("--template", required=True, help="Path to source template (.pptx).")
    parser.add_argument("--output", required=True, help="Output path for probe deck (.pptx).")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    result = build_probe_deck(Path(args.template), Path(args.output))
    print(f"Probe deck generated: {result}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
