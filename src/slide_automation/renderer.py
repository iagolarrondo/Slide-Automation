"""Deck rendering logic using python-pptx and template mapping."""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import re
import sys
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .template_registry.models import AGENDA_PLACEHOLDER_LIST, AGENDA_SLOAN_DONOR_TABLE
from .template_registry.sloan import SLIDE_TYPE_MAP as _DEFAULT_SLIDE_TYPE_MAP

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_R_ID_ATTR = f"{{{_R_NS}}}id"

# Donor-slide cloning: do not copy add-in / OLE / tags parts into output (PowerPoint repair issues).
_SKIP_DONOR_REL_TYPES = frozenset(
    {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags",
    }
)
_OLE_SHAPE_TYPES = frozenset(
    {
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
    }
)


def _xml_local_tag(tag: str) -> str:
    return tag.split("}")[-1] if "}" in tag else tag


def _element_tree_contains_ole(element: Any) -> bool:
    """True if subtree contains DrawingML OLE markup (add-ins, think-cell, etc.)."""
    for el in element.iter():
        t = _xml_local_tag(el.tag)
        if t in {"oleObj", "OLEObject"}:
            return True
    return False


def _should_skip_donor_shape_for_clone(shape: Any) -> bool:
    """Exclude OLE / add-in shapes; keep placeholders, pictures, tables, normal text."""
    name = (getattr(shape, "name", None) or "").lower()
    if "think-cell" in name:
        return True
    try:
        st = shape.shape_type
        if st in _OLE_SHAPE_TYPES:
            return True
    except Exception:
        pass
    try:
        if _element_tree_contains_ole(shape.element):
            return True
    except Exception:
        pass
    return False


CONTENT_MAIN_TITLE_TYPES = {
    "standard_1_block",
    "standard_2_block",
    "standard_3_block",
    "standard_2_block_big_left",
    "standard_2_block_big_right",
    "narrow_image_content",
    "wide_image_content",
    "wd_cover",
    "wd_section_intro",
    "wd_two_column",
    "wd_two_column_alt",
    "wd_three_column",
    "wd_four_column",
    "wd_mosaic_4",
    "wd_mosaic_6",
    "wd_mosaic_8",
    "wd_tile_4",
    "wd_tile_5",
    "wd_tile_6",
}


def _remove_slide_at(prs: Any, slide_index: int) -> None:
    """Remove slide at index from a presentation."""
    slide_ids = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_id = slide_ids[slide_index]
    prs.part.drop_rel(slide_id.rId)
    del slide_ids[slide_index]


def _clone_donor_slide(prs: Any, donor_slide: Any) -> Any:
    """Clone a donor slide into the presentation and return the new slide.

    python-pptx does not expose a first-class duplicate API, so we copy shape XML
    and non-notes relationships from the donor into a fresh slide with same layout.

    Add-in / OLE shapes (e.g. think-cell) and slide-level ``tags`` relationships are
    skipped so output decks do not embed those artifacts.
    """
    new_slide = prs.slides.add_slide(donor_slide.slide_layout)

    # Remove auto-generated layout shapes first.
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape.element)  # type: ignore[attr-defined]

    # Copy donor relationships (excluding OLE/tags) and create an old->new rId map.
    rid_map: Dict[str, str] = {}
    for rel in donor_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.reltype in _SKIP_DONOR_REL_TYPES:
            continue
        try:
            # python-pptx public-ish API for relationships is SlidePart.relate_to().
            new_rid = new_slide.part.relate_to(rel._target, rel.reltype)
        except Exception:
            continue
        if getattr(rel, "rId", None):
            rid_map[str(rel.rId)] = str(new_rid)

    # Copy donor shapes (excluding OLE / add-in artifacts).
    for shape in donor_slide.shapes:
        if _should_skip_donor_shape_for_clone(shape):
            continue
        new_el = deepcopy(shape.element)
        # Rewrite relationship attribute references (r:id, r:embed, r:link, etc.)
        # to match this slide part.
        for el in new_el.iter():
            # lxml-style expanded attribute keys are used in python-pptx oxml.
            for attr_name, attr_val in list(getattr(el, "attrib", {}).items()):
                if not isinstance(attr_name, str) or not attr_name.startswith(f"{{{_R_NS}}}"):
                    continue
                if isinstance(attr_val, str) and attr_val in rid_map:
                    el.set(attr_name, rid_map[attr_val])
        new_slide.shapes._spTree.insert_element_before(  # type: ignore[attr-defined]
            new_el,
            "p:extLst",
        )

    return new_slide


def _resolve_mapped_shape(slide: Any, key: str) -> Any:
    """Resolve a mapped target shape for a semantic key on the current slide.

    Mapping values can be:
      - int: placeholder idx
      - {"placeholder_idx": int}
      - {"shape_index": int}
    """
    mapping = slide._field_mapping  # type: ignore[attr-defined]
    target = mapping.get(key)

    placeholder_idx: Optional[int] = None
    shape_index: Optional[int] = None

    if isinstance(target, int):
        placeholder_idx = target
    elif isinstance(target, dict):
        if isinstance(target.get("placeholder_idx"), int):
            placeholder_idx = int(target["placeholder_idx"])
        if isinstance(target.get("shape_index"), int):
            shape_index = int(target["shape_index"])

    if placeholder_idx is not None:
        try:
            return slide.placeholders[placeholder_idx]
        except Exception:
            pass

    if shape_index is not None and 0 <= shape_index < len(slide.shapes):
        return slide.shapes[shape_index]

    return None


def _set_text(placeholder: Any, text: str) -> bool:
    """Set text in a placeholder safely.

    Returns:
        True when text was written, False otherwise.
    """
    if placeholder is None:
        return False
    try:
        placeholder.text = text
        return True
    except Exception:
        return False


def _normalize_text(value: Any) -> str:
    """Convert scalar/list values to placeholder text.

    Lists are rendered as plain lines so template bullet styles can apply.
    """
    def _strip_manual_bullet_prefix(text: str) -> str:
        # Remove common manual bullet prefixes; rely on PPT placeholder styling.
        return re.sub(r"^\s*(?:[-*•]\s+|\d+[\.\)]\s+)", "", text).strip()

    if isinstance(value, list):
        return "\n".join(_strip_manual_bullet_prefix(str(item)) for item in value)
    return _strip_manual_bullet_prefix(str(value))


def _clean_heading(text: str) -> str:
    """Light normalization for title-like fields: whitespace and stray edge punctuation."""
    text = " ".join(text.split()).strip()
    return text.strip(" ,;:-")


# Advisory character counts for stderr warnings only (editorial length is spec-owned).
_HEADING_ADVISORY_MAX_NONSPACE_CHARS = 86


def _nonspace_len(text: str) -> int:
    return len(re.sub(r"\s+", "", text))


def _block_title_max_chars(slide_type: str, field: str) -> int:
    # Warnings only. These reflect typical “block label” capacities by layout family.
    if slide_type in {"standard_3_block"}:
        return 18
    if slide_type in {"standard_2_block"} and field in {"left_block_title", "right_block_title"}:
        return 20
    if slide_type in {"standard_2_block_big_left", "standard_2_block_big_right"}:
        return 22
    if slide_type in {"narrow_image_content", "wide_image_content"}:
        return 22
    # Default single-block labels tolerate a bit more.
    return 24


def _main_title_max_chars(slide_type: str) -> int:
    """Advisory heading threshold (non-space chars). Warnings only."""
    return _HEADING_ADVISORY_MAX_NONSPACE_CHARS


def _warn_title_length(slide_type: str, field: str, text: str, limit: int) -> None:
    """Emit a non-fatal warning when a title-like field exceeds an advisory length."""
    n = _nonspace_len(text)
    if n <= limit:
        return
    print(
        f"[title-length] {slide_type}.{field}: {n} non-space chars (advisory max {limit}); "
        "rendering verbatim — shorten in the deck spec if layout overflow occurs",
        file=sys.stderr,
    )


def _prepare_main_title(text: str, slide_type: str) -> str:
    """Normalize whitespace; warn if likely long for layout; never shorten semantically."""
    prepared = _clean_heading(text)
    _warn_title_length(slide_type, "title", prepared, _main_title_max_chars(slide_type))
    return prepared


def _prepare_block_title(text: str, slide_type: str, field: str) -> str:
    """Normalize whitespace; warn if likely long for layout; never shorten semantically."""
    prepared = _clean_heading(text)
    _warn_title_length(slide_type, field, prepared, _block_title_max_chars(slide_type, field))
    return prepared


def _set_title(slide: Any, text: str, slide_type: str) -> bool:
    """Set slide title using mapped placeholder first, then fallback title shape."""
    if slide_type in CONTENT_MAIN_TITLE_TYPES:
        text = _prepare_main_title(text, slide_type)
    else:
        text = _clean_heading(text)
    ph = _resolve_mapped_shape(slide, "title")
    if _set_text(ph, text):
        return True

    # Fallback: native title shape when available.
    if getattr(slide.shapes, "title", None) is not None:
        try:
            slide.shapes.title.text = text
            return True
        except Exception:
            return False

    return False


def _set_body_or_placeholder(slide: Any, key: str, value: Any) -> bool:
    """Set text to mapped placeholder key or fallback to a body placeholder."""
    mapping = slide._field_mapping  # type: ignore[attr-defined]
    target = mapping.get(key)
    ph = _resolve_mapped_shape(slide, key)
    text = _normalize_text(value)
    if _set_text(ph, text):
        return True

    # If an explicit mapping target is unavailable, do not fallback into unrelated
    # placeholders (prevents accidental content overwrite).
    if target is not None and ph is None:
        return False

    # Fallback to first non-title placeholder with text frame.
    for candidate in slide.placeholders:
        try:
            if candidate == slide.shapes.title:
                continue
        except Exception:
            pass

        if hasattr(candidate, "text_frame"):
            try:
                candidate.text = text
                return True
            except Exception:
                continue
    return False


def _set_footer_and_slide_number(
    slide: Any,
    footer_value: Any,
    slide_number: int,
) -> None:
    """Write footer and slide number only into native materialized placeholders.

    python-pptx often does not materialize FOOTER / SLIDE_NUMBER placeholders on
    slides created from layouts. In that case we skip silently (no ad hoc shapes).
    """
    mapping = slide._field_mapping  # type: ignore[attr-defined]
    footer_target = mapping.get("footer")
    footer_text = _normalize_text(footer_value).strip()

    footer_ph = _resolve_mapped_shape(slide, "footer")

    # Determine slide-number placeholder by scanning materialized placeholders.
    slide_number_idx: Optional[int] = None
    for ph in slide.placeholders:
        try:
            if str(ph.placeholder_format.type).startswith("SLIDE_NUMBER"):
                slide_number_idx = ph.placeholder_format.idx
                break
        except Exception:
            continue

    # Footer policy: never retain donor text like "Footer".
    # - If footer provided: set it.
    # - Else: clear placeholder text entirely.
    if footer_ph is not None:
        try:
            ph_type = str(footer_ph.placeholder_format.type)
        except Exception:
            ph_type = ""
        if ph_type.startswith("FOOTER"):
            _set_text(footer_ph, footer_text)

    if isinstance(slide_number_idx, int):
        try:
            slide_number_ph = slide.placeholders[slide_number_idx]
        except Exception:
            slide_number_ph = None
        if slide_number_ph is not None:
            try:
                ph_type = str(slide_number_ph.placeholder_format.type)
            except Exception:
                ph_type = ""
            if ph_type.startswith("SLIDE_NUMBER"):
                _set_text(slide_number_ph, str(slide_number))


def _set_image_on_placeholder(slide: Any, key: str, image_path: Any) -> bool:
    """Insert image into mapped picture placeholder (or geometry fallback)."""
    path = Path(str(image_path)) if image_path is not None else None
    if path is None or not path.exists():
        return False

    target_shape = _resolve_mapped_shape(slide, key)
    if target_shape is None:
        return False

    try:
        target_shape.insert_picture(str(path))
        return True
    except Exception:
        pass

    # Fallback to placeholder geometry if direct insert is unavailable.
    try:
        slide.shapes.add_picture(
            str(path),
            target_shape.left,
            target_shape.top,
            target_shape.width,
            target_shape.height,
        )
        return True
    except Exception:
        return False


_BAR_CHARS = frozenset("|│\u2502")


def _strip_duplicate_leading_bars(text: str) -> str:
    t = text.strip()
    while t and t[0] in _BAR_CHARS:
        t = t[1:].lstrip()
    return t


def _detect_agenda_label_column(table: Any) -> int:
    """Prefer the donor column that holds the │ section labels (often col 1 when col 0 is blank)."""
    ncol = len(table.columns)
    if ncol == 1:
        return 0
    try:
        left = table.cell(0, 0).text_frame.text.strip()
        right = table.cell(0, 1).text_frame.text.strip()
        if not left and right:
            return 1
        if left and not right:
            return 0
    except Exception:
        pass
    return ncol - 1


def _donor_agenda_separator_and_two_runs(cell: Any) -> tuple[str, bool]:
    """Read leading bar from donor label cell; detect multi-run │ + title pattern."""
    try:
        p = cell.text_frame.paragraphs[0]
    except Exception:
        return "\u2502", False
    runs = p.runs
    if len(runs) >= 2:
        head = runs[0].text.strip()
        if head and all(c in _BAR_CHARS for c in head):
            return runs[0].text, True
    full = "".join(r.text for r in runs)
    m = re.match(r"^([\|│\u2502])\s*", full.strip())
    if m:
        return m.group(1), False
    return "\u2502", False


def _best_agenda_template_row_idx(table: Any, label_col: int) -> int:
    """Pick a donor row whose label cell has ≥2 runs so clones inherit sep + title formatting."""
    for i in range(len(table.rows)):
        try:
            cell = table.cell(i, label_col)
            if len(cell.text_frame.paragraphs[0].runs) >= 2:
                return i
        except Exception:
            continue
    return 0


def _replace_cell_text_without_clear(cell: Any, text: str) -> None:
    """Replace visible text without calling TextFrame.clear() (preserves a:rPr formatting)."""
    tf = cell.text_frame
    if not tf.paragraphs:
        tf.text = text
        return
    p = tf.paragraphs[0]
    rs = p.runs
    if rs:
        rs[0].text = text
        for r in rs[1:]:
            r.text = ""
    else:
        p.text = text


def _clear_cell_visible_text(cell: Any) -> None:
    """Clear cell text while keeping paragraph/run scaffolding when present."""
    try:
        tf = cell.text_frame
        if not tf.paragraphs:
            return
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = ""
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = ""
    except Exception:
        pass


def _apply_agenda_label_cell(cell: Any, sep: str, label: str, prefer_two_runs: bool) -> None:
    """Write section label with exactly one leading bar (donor-style │ or ASCII |)."""
    body = _strip_duplicate_leading_bars(label)
    p = cell.text_frame.paragraphs[0]
    runs = p.runs
    if prefer_two_runs and len(runs) >= 2:
        runs[0].text = sep
        runs[1].text = body
        for r in runs[2:]:
            r.text = ""
        return
    if runs:
        runs[0].text = sep + body
        for r in runs[1:]:
            r.text = ""
    else:
        p.text = sep + body


def _duplicate_picture_shape_on_slide(slide: Any, template_pic: Any) -> None:
    """Clone a picture shape on the same slide (SVG/PNG); maps r:embed to new relationships.

    Donor MIT logos often use SVG ``asvg:svgBlip``; ``add_picture``/``.image.blob`` may not apply.
    """
    new_el = deepcopy(template_pic.element)
    rid_map: Dict[str, str] = {}
    slide_part = slide.part
    for el in new_el.iter():
        for attr_name, attr_val in list(el.attrib.items()):
            if not isinstance(attr_name, str) or not attr_name.startswith(f"{{{_R_NS}}}"):
                continue
            if not isinstance(attr_val, str) or not attr_val.startswith("rId"):
                continue
            if attr_val in rid_map:
                continue
            try:
                rel = slide_part.rels[attr_val]
                new_rid = slide_part.relate_to(rel._target, rel.reltype)
                rid_map[attr_val] = str(new_rid)
            except Exception:
                continue
    for el in new_el.iter():
        for attr_name, attr_val in list(el.attrib.items()):
            if (
                isinstance(attr_name, str)
                and attr_name.startswith(f"{{{_R_NS}}}")
                and attr_val in rid_map
            ):
                el.set(attr_name, rid_map[attr_val])
    slide.shapes._spTree.insert_element_before(  # type: ignore[attr-defined]
        new_el,
        "p:extLst",
    )


def _agenda_slide_pictures_by_row(slide: Any, table_shape: Any) -> List[Any]:
    """MIT logos: picture shapes on the agenda slide (one per row), excluding the table itself."""
    pics: List[Any] = []
    for sh in slide.shapes:
        if sh is table_shape:
            continue
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        pics.append(sh)
    pics.sort(key=lambda s: int(getattr(s, "top", 0) or 0))
    return pics


def _resize_agenda_table_rows(table: Any, n: int, template_tr_idx: int) -> None:
    """Add or remove ``a:tr`` elements at the XML level; update graphic frame height."""
    tbl = table._tbl
    template_tr_idx = max(0, min(template_tr_idx, len(tbl.tr_lst) - 1))
    while len(tbl.tr_lst) > n:
        tr = tbl.tr_lst[-1]
        tr.getparent().remove(tr)
    while len(tbl.tr_lst) < n:
        new_tr = deepcopy(tbl.tr_lst[template_tr_idx])
        tbl.append(new_tr)
    table.notify_height_changed()


def _agenda_vertically_center_block(
    slide_height: int, table_shape: Any, logo_shapes: List[Any]
) -> None:
    """Move table + logos together so the block is centered on the slide."""
    tops: List[int] = [int(table_shape.top)]
    bottoms: List[int] = [int(table_shape.top) + int(table_shape.height)]
    for logo in logo_shapes:
        tops.append(int(logo.top))
        bottoms.append(int(logo.top) + int(logo.height))
    block_top = min(tops)
    block_bottom = max(bottoms)
    cy = (block_top + block_bottom) // 2
    target = int(slide_height) // 2
    delta = target - cy
    if delta == 0:
        return
    table_shape.top = int(table_shape.top) + delta
    for logo in logo_shapes:
        logo.top = int(logo.top) + delta


def _render_agenda_table(pres: Any, slide: Any, agenda_items: Any) -> bool:
    """Render agenda items using donor table + row-aligned MIT logos (donor-deck path).

    Preserves donor cell typography by avoiding ``TextFrame.clear()`` and by editing runs
    in place. Coordinates table row count with logo shapes and re-centers the composition.
    """
    items = agenda_items if isinstance(agenda_items, list) else []
    items_text = [str(item) for item in items]

    target_shape = _resolve_mapped_shape(slide, "agenda_items")

    def _write_plain_cell(cell: Any, text: str, align: Any) -> None:
        text_frame = cell.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    if target_shape is None or not getattr(target_shape, "has_table", False):
        if target_shape is not None:
            try:
                rows_ct = max(1, len(items_text))
                cols = 2
                shape = slide.shapes.add_table(
                    rows_ct,
                    cols,
                    target_shape.left,
                    target_shape.top,
                    target_shape.width,
                    target_shape.height,
                )
                table = shape.table
                table.columns[0].width = int(target_shape.width * 0.16)
                table.columns[1].width = int(target_shape.width * 0.84)
                for row_idx in range(rows_ct):
                    title_text = items_text[row_idx] if row_idx < len(items_text) else ""
                    _write_plain_cell(table.cell(row_idx, 0), "", PP_ALIGN.CENTER)
                    _write_plain_cell(table.cell(row_idx, 1), title_text, PP_ALIGN.LEFT)
                return True
            except Exception:
                pass
        return _set_body_or_placeholder(slide, "agenda_items", items_text)

    table = target_shape.table
    if len(table.rows) == 0 or not items_text:
        return False

    label_col = _detect_agenda_label_column(table)
    num_col = 0 if label_col != 0 else None
    sep, prefer_two_runs_hint = _donor_agenda_separator_and_two_runs(table.cell(0, label_col))
    template_idx = _best_agenda_template_row_idx(table, label_col)
    n_rows_donor = len(table.rows)
    clone_src_idx = min(template_idx, max(0, n_rows_donor - 1))
    try:
        tpl_cell = table.cell(clone_src_idx, label_col)
        prefer_two_runs = len(tpl_cell.text_frame.paragraphs[0].runs) >= 2
    except Exception:
        prefer_two_runs = prefer_two_runs_hint

    n = len(items_text)
    _resize_agenda_table_rows(table, n, clone_src_idx)

    for i in range(n):
        if num_col is not None:
            _clear_cell_visible_text(table.cell(i, num_col))
        _apply_agenda_label_cell(
            table.cell(i, label_col),
            sep,
            items_text[i],
            prefer_two_runs,
        )

    slide_h = getattr(pres, "slide_height", None)
    if slide_h is None:
        slide_h = 6858000
    slide_h = int(slide_h)

    logos = _agenda_slide_pictures_by_row(slide, target_shape)
    if not logos:
        _agenda_vertically_center_block(slide_h, target_shape, [])
        return True

    ref_left = int(logos[0].left)
    ref_w = int(logos[0].width)
    ref_h = int(logos[0].height)

    if len(logos) > n:
        for sh in logos[n:]:
            try:
                sh._element.getparent().remove(sh._element)  # type: ignore[attr-defined]
            except Exception:
                pass
        logos = logos[:n]

    if logos and n > 0:
        template_logo = logos[0]
        for _ in range(max(0, n - len(logos))):
            _duplicate_picture_shape_on_slide(slide, template_logo)
        logos = _agenda_slide_pictures_by_row(slide, target_shape)

    for i, logo in enumerate(logos[:n]):
        row_top = int(target_shape.top) + sum(
            int(table.rows[j].height) for j in range(i)
        )
        rh = int(table.rows[i].height)
        logo.left = ref_left
        logo.width = ref_w
        logo.height = ref_h
        logo.top = int(row_top + (rh - ref_h) // 2)

    _agenda_vertically_center_block(slide_h, target_shape, logos[:n])

    return True


def render_deck(
    template_path: str | Path,
    payload: Dict[str, Any],
    output_path: str | Path,
    *,
    slide_type_map: Dict[str, Dict[str, Any]] | None = None,
    agenda_render_mode: str | None = None,
) -> Path:
    """Render a PPTX file from deck payload and template.

    Args:
        template_path: Existing .potx or .pptx template path.
        payload: Parsed and validated deck payload.
        output_path: Destination .pptx path.
        slide_type_map: Per-slide donor mapping; defaults to the Sloan bundled map.
        agenda_render_mode: Used only for Sloan ``agenda`` slides (table vs placeholder list).
            WD agendas use distinct ``wd_agenda_*`` slide types and do not use this path.

    Returns:
        Resolved output path.

    Raises:
        FileNotFoundError: If template is missing.
        ValueError: If slide type mapping is invalid.
    """
    type_map = slide_type_map if slide_type_map is not None else _DEFAULT_SLIDE_TYPE_MAP
    agenda_mode = agenda_render_mode or AGENDA_SLOAN_DONOR_TABLE
    template = Path(template_path)
    if not template.exists():
        raise FileNotFoundError(f"Template file not found: {template}")

    prs = Presentation(str(template))
    donor_slides = list(prs.slides)
    donor_count = len(donor_slides)

    for slide_number, slide_data in enumerate(payload.get("slides", []), start=1):
        slide_type = slide_data.get("type")
        if slide_type not in type_map:
            raise ValueError(f"Unsupported slide type in mapping: {slide_type!r}")

        config = type_map[slide_type]
        donor_slide_number = config.get("donor_slide_number")
        expected_layout_name = str(config.get("expected_layout_name", ""))
        placeholder_map = config.get("placeholders", {})

        if not isinstance(donor_slide_number, int):
            raise ValueError(
                f"donor_slide_number for slide type '{slide_type}' must be an integer."
            )

        donor_index = donor_slide_number - 1
        if donor_index < 0 or donor_index >= donor_count:
            raise ValueError(
                f"donor_slide_number {donor_slide_number} for '{slide_type}' is out of range "
                f"(template has {donor_count} donor slides)."
            )

        donor_slide = donor_slides[donor_index]
        actual_layout_name = str(donor_slide.slide_layout.name)
        if expected_layout_name and actual_layout_name != expected_layout_name:
            raise ValueError(
                f"donor slide {donor_slide_number} for '{slide_type}' has layout "
                f"{actual_layout_name!r}, expected {expected_layout_name!r}."
            )

        slide = _clone_donor_slide(prs, donor_slide)
        slide._field_mapping = placeholder_map  # type: ignore[attr-defined]

        if slide_type.startswith("wd_"):
            from .template_registry.wd_render import apply_wd_slide

            apply_wd_slide(prs, slide, slide_type, slide_data, slide_number, config)
            # WD path owns its own editable-target policy and footer writes.
            # Keep slide-number placeholders untouched by generic footer logic.
            continue

        # Agenda is a special layout: title text is baked into template.
        if slide_type != "agenda":
            title = str(slide_data.get("title", ""))
            if title:
                _set_title(slide, title, slide_type)

        if slide_type == "cover":
            if "subtitle" in slide_data:
                _set_body_or_placeholder(slide, "subtitle", slide_data.get("subtitle", ""))
            if "date" in slide_data:
                _set_body_or_placeholder(slide, "date", slide_data.get("date", ""))

        elif slide_type == "agenda":
            if "agenda_items" in slide_data:
                items = slide_data.get("agenda_items", [])
                if agenda_mode == AGENDA_SLOAN_DONOR_TABLE:
                    _render_agenda_table(prs, slide, items)
                elif agenda_mode == AGENDA_PLACEHOLDER_LIST:
                    _set_body_or_placeholder(slide, "agenda_items", items)
                else:
                    raise ValueError(f"Unknown agenda_render_mode: {agenda_mode!r}")

        elif slide_type == "divider":
            # Divider contract:
            # - title -> main divider heading (mapped "title" placeholder)
            # - section_number (optional) -> numeric box, zero-padded 2 digits
            # - section_title (legacy/back-compat) may contain "Section 2"; we extract the number.
            raw_section_number = slide_data.get("section_number", None)
            if raw_section_number is None:
                legacy = str(slide_data.get("section_title", "")).strip()
                m = re.search(r"(\d+)", legacy)
                raw_section_number = int(m.group(1)) if m else None
            if raw_section_number is not None:
                try:
                    n = int(raw_section_number)
                except Exception:
                    n = None
                if isinstance(n, int) and n >= 0:
                    _set_body_or_placeholder(slide, "section_number", f"{n:02d}")

        elif slide_type == "standard_1_block":
            for key in ("section_title", "block_title", "block_body"):
                if key in slide_data:
                    if key == "block_title":
                        original = str(slide_data.get(key, ""))
                        prepared = _prepare_block_title(original, slide_type, key)
                        _set_body_or_placeholder(slide, key, prepared)
                    else:
                        _set_body_or_placeholder(slide, key, slide_data.get(key, ""))
            _set_footer_and_slide_number(slide, slide_data.get("footer", ""), slide_number)

        elif slide_type == "standard_2_block":
            for key in (
                "section_title",
                "left_block_title",
                "left_block_body",
                "right_block_title",
                "right_block_body",
            ):
                if key in slide_data:
                    if key in {"left_block_title", "right_block_title"}:
                        original = str(slide_data.get(key, ""))
                        prepared = _prepare_block_title(original, slide_type, key)
                        _set_body_or_placeholder(slide, key, prepared)
                    else:
                        _set_body_or_placeholder(slide, key, slide_data.get(key, ""))
            _set_footer_and_slide_number(slide, slide_data.get("footer", ""), slide_number)

        elif slide_type == "standard_3_block":
            for key in (
                "section_title",
                "block_1_title",
                "block_1_body",
                "block_2_title",
                "block_2_body",
                "block_3_title",
                "block_3_body",
            ):
                if key in slide_data:
                    if key in {"block_1_title", "block_2_title", "block_3_title"}:
                        original = str(slide_data.get(key, ""))
                        prepared = _prepare_block_title(original, slide_type, key)
                        _set_body_or_placeholder(slide, key, prepared)
                    else:
                        _set_body_or_placeholder(slide, key, slide_data.get(key, ""))
            _set_footer_and_slide_number(slide, slide_data.get("footer", ""), slide_number)

        elif slide_type in {"standard_2_block_big_left", "standard_2_block_big_right"}:
            for key in (
                "section_title",
                "dominant_block_title",
                "dominant_block_body",
                "secondary_block_title",
                "secondary_block_body",
            ):
                if key in slide_data:
                    if key in {"dominant_block_title", "secondary_block_title"}:
                        original = str(slide_data.get(key, ""))
                        prepared = _prepare_block_title(original, slide_type, key)
                        _set_body_or_placeholder(slide, key, prepared)
                    else:
                        _set_body_or_placeholder(slide, key, slide_data.get(key, ""))
            _set_footer_and_slide_number(slide, slide_data.get("footer", ""), slide_number)

        elif slide_type in {"narrow_image_content", "wide_image_content"}:
            for key in ("section_title", "content_block_title", "content_block_body"):
                if key in slide_data:
                    if key == "content_block_title":
                        original = str(slide_data.get(key, ""))
                        prepared = _prepare_block_title(original, slide_type, key)
                        _set_body_or_placeholder(slide, key, prepared)
                    else:
                        _set_body_or_placeholder(slide, key, slide_data.get(key, ""))
            if "image" in slide_data:
                _set_image_on_placeholder(slide, "image", slide_data.get("image"))
            _set_footer_and_slide_number(slide, slide_data.get("footer", ""), slide_number)

    # Remove canonical donor slides from output.
    for _ in range(donor_count):
        _remove_slide_at(prs, 0)

    destination = Path(output_path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destination))
    return destination.resolve()
