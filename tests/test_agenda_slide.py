"""Regression tests for donor-deck agenda slide (table + row logos)."""

from __future__ import annotations

import re
import tempfile
import unittest
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

from slide_automation.renderer import render_deck


def _agenda_slide_xml(z: zipfile.ZipFile) -> bytes:
    names = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    names = [n for n in names if "_rels" not in n]
    names.sort()
    return z.read(names[0])


def _table_cell_texts_from_slide_xml(xml_bytes: bytes) -> list[list[str]]:
    root = ET.fromstring(xml_bytes)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    rows: list[list[str]] = []
    for tr in root.findall(".//a:tbl/a:tr", ns):
        row_cells: list[str] = []
        for tc in tr.findall("a:tc", ns):
            texts: list[str] = []
            for t in tc.findall(".//a:t", ns):
                if t.text:
                    texts.append(t.text)
            row_cells.append("".join(texts))
        rows.append(row_cells)
    return rows


def _count_leading_bars(cell_text: str) -> int:
    n = 0
    for ch in cell_text:
        if ch in "|│\u2502":
            n += 1
        elif ch.isspace():
            continue
        else:
            break
    return n


def _block_center_vs_slide_mid(prs: Presentation, slide) -> tuple[int, int]:
    gf = next(s for s in slide.shapes if getattr(s, "has_table", False))
    table = gf
    pics = [s for s in slide.shapes if s.shape_type == 13]
    tops = [table.top] + [p.top for p in pics]
    bottoms = [table.top + table.height] + [p.top + p.height for p in pics]
    cy = (min(tops) + max(bottoms)) // 2
    mid = int(prs.slide_height or 0) // 2
    return cy, mid


class TestAgendaDonorSlide(unittest.TestCase):
    def setUp(self) -> None:
        self.repo = Path(__file__).resolve().parents[1]
        self.donor = self.repo / "templates" / "Sloan_Donor_Deck.pptx"

    def _render(self, items: list[str]) -> Path:
        payload = {"deck_title": "T", "slides": [{"type": "agenda", "agenda_items": items}]}
        tmp = Path(tempfile.mkdtemp()) / "out.pptx"
        render_deck(self.donor, payload, tmp)
        return tmp

    def test_no_section_numbers_in_table(self) -> None:
        out = self._render(["First", "Second", "Third"])
        prs = Presentation(str(out))
        slide = prs.slides[0]
        gf = next(s for s in slide.shapes if getattr(s, "has_table", False))
        table = gf.table
        for i in range(len(table.rows)):
            left = table.cell(i, 0).text_frame.text.strip()
            self.assertEqual(left, "", msg="agenda must not use the narrow column for numbering")
            lab = table.cell(i, 1).text_frame.text
            self.assertIsNone(re.match(r"^\s*\d+", lab), msg=f"unexpected digit-leading label: {lab!r}")

    def test_exactly_one_bar_prefix_per_row(self) -> None:
        out = self._render(["Alpha", "Beta"])
        xml = _agenda_slide_xml(zipfile.ZipFile(out, "r"))
        rows = _table_cell_texts_from_slide_xml(xml)
        for row in rows:
            self.assertGreaterEqual(len(row), 2)
            label = row[1]
            self.assertGreaterEqual(len(label), 2)
            self.assertEqual(_count_leading_bars(label), 1, msg=f"label must have one bar prefix: {label!r}")

    def test_fewer_rows_drops_extra_logos(self) -> None:
        out = self._render(["Only", "Two"])
        prs = Presentation(str(out))
        slide = prs.slides[0]
        pics = [s for s in slide.shapes if s.shape_type == 13]
        self.assertEqual(len(pics), 2)

    def test_more_rows_duplicates_logos(self) -> None:
        out = self._render([f"Row {i}" for i in range(7)])
        prs = Presentation(str(out))
        slide = prs.slides[0]
        gf = next(s for s in slide.shapes if getattr(s, "has_table", False))
        table = gf.table
        pics = [s for s in slide.shapes if s.shape_type == 13]
        self.assertEqual(len(table.rows), 7)
        self.assertEqual(len(pics), 7)

    def test_block_vertically_centered(self) -> None:
        for n in (3, 7):
            out = self._render([f"S{i}" for i in range(n)])
            prs = Presentation(str(out))
            cy, mid = _block_center_vs_slide_mid(prs, prs.slides[0])
            self.assertEqual(cy, mid, msg=f"N={n}: agenda block should be vertically centered")

    def test_label_cell_preserves_run_properties(self) -> None:
        """Donor formatting uses a:rPr (e.g. sz); we must not strip it via TextFrame.clear()."""
        out = self._render(["My Section"])
        xml = _agenda_slide_xml(zipfile.ZipFile(out, "r"))
        root = ET.fromstring(xml)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        rprs = root.findall(".//a:tbl//a:tc//a:r/a:rPr", ns)
        self.assertTrue(rprs, msg="expected at least one a:rPr on agenda table text runs")


if __name__ == "__main__":
    unittest.main()
