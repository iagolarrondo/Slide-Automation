"""Regression tests for WD donor fidelity and safe editable-target writes."""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slide_automation.renderer import render_deck
from slide_automation.template_registry.wd import WD_SLIDE_TYPE_MAP


def _render_wd(payload: dict) -> Presentation:
    repo = Path(__file__).resolve().parents[1]
    donor = repo / "templates" / "WD_Template_Donor.pptx"
    out = Path(tempfile.mkdtemp()) / "wd_out.pptx"
    render_deck(
        donor,
        payload,
        out,
        slide_type_map=WD_SLIDE_TYPE_MAP,
    )
    return Presentation(str(out))


def _first_run_style(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if not p.runs:
        return None
    r = p.runs[0]
    c = r.font.color
    return (
        str(getattr(r.font, "bold", None)),
        str(getattr(r.font, "italic", None)),
        str(getattr(r.font, "size", None)),
        str(getattr(c, "type", None)),
        str(getattr(c, "rgb", None)),
        str(getattr(c, "theme_color", None)),
    )


def _all_texts_recursive(slide):
    out = []

    def walk(shape):
        if hasattr(shape, "text_frame"):
            out.append(shape.text_frame.text or "")
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                walk(sub)

    for sh in slide.shapes:
        walk(sh)
    return out


class TestWdRenderingFidelity(unittest.TestCase):
    def setUp(self) -> None:
        repo = Path(__file__).resolve().parents[1]
        self.donor = Presentation(str(repo / "templates" / "WD_Template_Donor.pptx"))

    def test_cover_preserves_title_style(self) -> None:
        donor_title = self.donor.slides[0].placeholders[0]
        donor_style = _first_run_style(donor_title)
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_cover",
                        "title": "New White Title",
                        "subtitle": "Sub",
                        "presenter": "Presenter",
                        "date": "Date",
                    }
                ],
            }
        )
        out_title = prs.slides[0].placeholders[0]
        self.assertEqual(out_title.text_frame.text.strip(), "New White Title")
        self.assertEqual(_first_run_style(out_title), donor_style)

    def test_divider_does_not_change_numeral_boxes(self) -> None:
        donor = self.donor.slides[5]
        before = set()
        for sh in donor.shapes:
            if hasattr(sh, "text_frame"):
                txt = (sh.text_frame.text or "").strip()
                if txt.isdigit() and len(txt) <= 2:
                    before.add(txt)
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [{"type": "wd_divider", "title": "Only Title"}],
            }
        )
        out = prs.slides[0]
        after = set()
        for sh in out.shapes:
            if hasattr(sh, "text_frame"):
                txt = (sh.text_frame.text or "").strip()
                if txt.isdigit() and len(txt) <= 2:
                    after.add(txt)
        self.assertEqual(after, before)
        # Ensure title is written into the divider title paragraph (not numeral paragraph).
        divider_text = next(
            sh.text_frame.text
            for sh in out.shapes
            if hasattr(sh, "text_frame") and "Only Title" in (sh.text_frame.text or "")
        )
        self.assertTrue(divider_text.startswith("01"))
        self.assertIn("Only Title", divider_text)

    def test_wd_divider_uses_explicit_section_number_zero_padded(self) -> None:
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [{"type": "wd_divider", "title": "D", "section_number": 7}],
            }
        )
        txt = next(
            sh.text_frame.text
            for sh in prs.slides[0].shapes
            if hasattr(sh, "text_frame") and "D" in (sh.text_frame.text or "")
        )
        self.assertTrue(txt.startswith("07"))

    def test_wd_divider_falls_back_to_wd_divider_ordinal(self) -> None:
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {"type": "wd_divider", "title": "First"},
                    {"type": "wd_section_intro", "title": "mid", "section_title": "mid"},
                    {"type": "wd_divider", "title": "Second"},
                    {"type": "wd_divider", "title": "Third"},
                ],
            }
        )
        divider_texts = []
        for slide in prs.slides:
            for sh in slide.shapes:
                if hasattr(sh, "text_frame"):
                    t = sh.text_frame.text or ""
                    if "First" in t or "Second" in t or "Third" in t:
                        divider_texts.append(t)
        self.assertEqual(len(divider_texts), 3)
        self.assertTrue(divider_texts[0].startswith("01"))
        self.assertTrue(divider_texts[1].startswith("02"))
        self.assertTrue(divider_texts[2].startswith("03"))

    def test_agenda_populates_all_rows_and_clears_unused(self) -> None:
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_agenda_6",
                        "title": "Agenda",
                        "agenda_items": ["A", "B", "C", "D"],
                    }
                ],
            }
        )
        slide = prs.slides[0]
        texts = []
        for sh in slide.shapes:
            if hasattr(sh, "text_frame"):
                texts.append((sh.text_frame.text or "").strip())
        self.assertIn("A", texts)
        self.assertIn("B", texts)
        self.assertIn("C", texts)
        self.assertIn("D", texts)
        self.assertNotIn("Section Title", texts)

    def test_footer_written_only_to_footer_target(self) -> None:
        footer_text = "FOOTER_SENTINEL_123"
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_two_column",
                        "title": "Title",
                        "section_title": "Section",
                        "left_block_title": "L",
                        "left_block_body": "LB",
                        "right_block_title": "R",
                        "right_block_body": "RB",
                        "footer": footer_text,
                    }
                ],
            }
        )
        slide = prs.slides[0]
        hits = 0
        for sh in slide.shapes:
            if hasattr(sh, "text_frame") and footer_text in (sh.text_frame.text or ""):
                hits += 1
        self.assertEqual(hits, 1)

    def test_content_write_preserves_run_style_metadata(self) -> None:
        donor_slide = self.donor.slides[7]  # donor #8, wd_two_column
        donor_left = donor_slide.shapes[1]
        donor_style = _first_run_style(donor_left)
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_two_column",
                        "title": "Title",
                        "section_title": "Section",
                        "left_block_title": "Left",
                        "left_block_body": ["Line 1", "Line 2"],
                        "right_block_title": "Right",
                        "right_block_body": "Body",
                    }
                ],
            }
        )
        out_left = None
        for sh in prs.slides[0].shapes:
            if hasattr(sh, "text_frame") and "Left" in (sh.text_frame.text or ""):
                out_left = sh
                break
        self.assertIsNotNone(out_left)
        self.assertEqual(_first_run_style(out_left), donor_style)

    def test_two_column_alt_populates_body(self) -> None:
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_two_column_alt",
                        "title": "Title",
                        "section_title": "Section",
                        "left_block_title": "LTitle",
                        "left_block_body": "LBody",
                        "right_block_title": "RTitle",
                        "right_block_body": "RBody",
                    }
                ],
            }
        )
        texts = [
            sh.text_frame.text
            for sh in prs.slides[0].shapes
            if hasattr(sh, "text_frame")
        ]
        self.assertTrue(any("LTitle" in t and "LBody" in t for t in texts))
        self.assertTrue(any("RTitle" in t and "RBody" in t for t in texts))

    def test_three_and_four_column_populate_bodies(self) -> None:
        prs3 = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_three_column",
                        "title": "T",
                        "section_title": "S",
                        "block_1_title": "A",
                        "block_1_body": "A body",
                        "block_2_title": "B",
                        "block_2_body": "B body",
                        "block_3_title": "C",
                        "block_3_body": "C body",
                    }
                ],
            }
        )
        t3 = [sh.text_frame.text for sh in prs3.slides[0].shapes if hasattr(sh, "text_frame")]
        self.assertTrue(any("A body" in t for t in t3))
        self.assertTrue(any("B body" in t for t in t3))
        self.assertTrue(any("C body" in t for t in t3))

        prs4 = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_four_column",
                        "title": "T",
                        "section_title": "S",
                        "block_1_title": "A",
                        "block_1_body": "A body",
                        "block_2_title": "B",
                        "block_2_body": "B body",
                        "block_3_title": "C",
                        "block_3_body": "C body",
                        "block_4_title": "D",
                        "block_4_body": "D body",
                    }
                ],
            }
        )
        t4 = [sh.text_frame.text for sh in prs4.slides[0].shapes if hasattr(sh, "text_frame")]
        self.assertTrue(any("A body" in t for t in t4))
        self.assertTrue(any("B body" in t for t in t4))
        self.assertTrue(any("C body" in t for t in t4))
        self.assertTrue(any("D body" in t for t in t4))

    def test_one_block_grouped_semantics(self) -> None:
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_one_block_grouped",
                        "title": "Grouped Main",
                        "section_title": "Grouped Strap",
                        "block_title": "Grouped Cap",
                        "block_body": "Grouped Body",
                    }
                ],
            }
        )
        all_text = "\n".join(_all_texts_recursive(prs.slides[0]))
        self.assertIn("Grouped Cap", all_text)
        self.assertIn("Grouped Body", all_text)

    def test_one_block_placeholder_combines_title_body(self) -> None:
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_one_block_placeholder",
                        "title": "Placeholder Main",
                        "section_title": "Placeholder Strap",
                        "block_title": "Placeholder Cap",
                        "block_body": ["Line A", "Line B"],
                    }
                ],
            }
        )
        texts = [sh.text_frame.text for sh in prs.slides[0].shapes if hasattr(sh, "text_frame")]
        self.assertTrue(any("Placeholder Cap" in t and "Line A" in t and "Line B" in t for t in texts))

    def test_slide7_box_title_body_semantics(self) -> None:
        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_section_intro",
                        "title": "S7 Main",
                        "section_title": "S7 Strap",
                        "box_1_title": "LEFT CAP",
                        "box_1_body": "LEFT BODY",
                        "box_2_title": "RIGHT CAP",
                        "box_2_body": "RIGHT BODY",
                    }
                ],
            }
        )
        slide = prs.slides[0]
        texts = _all_texts_recursive(slide)
        all_text = "\n".join(texts)
        self.assertIn("LEFT CAP", all_text)
        self.assertIn("LEFT BODY", all_text)
        self.assertIn("RIGHT CAP", all_text)
        self.assertIn("RIGHT BODY", all_text)
        # Ensure donor box placeholders were replaced.
        self.assertNotIn("Title 1", all_text)
        self.assertNotIn("Title 2", all_text)
        self.assertNotIn("Text 1", all_text)
        self.assertNotIn("Text 2", all_text)

    def test_full_catalog_has_no_unreplaced_placeholder_text(self) -> None:
        repo = Path(__file__).resolve().parents[1]
        payload_path = repo / "examples" / "example_deck_wd_full_catalog.json"
        donor = repo / "templates" / "WD_Template_Donor.pptx"
        out = Path(tempfile.mkdtemp()) / "wd_full_catalog.pptx"
        import json

        payload = json.loads(payload_path.read_text(encoding="utf-8"))
        render_deck(donor, payload, out, slide_type_map=WD_SLIDE_TYPE_MAP)
        prs = Presentation(str(out))
        forbidden = {"Heading", "Section Title", "Title 1", "Title 2", "Text 1", "Text 2"}
        for slide in prs.slides:
            for txt in _all_texts_recursive(slide):
                txt = txt.strip()
                self.assertNotIn(txt, forbidden)

    def test_mosaic_title_body_semantics_and_style(self) -> None:
        donor_slide = self.donor.slides[16]  # donor #17 (wd_mosaic_4)
        donor_tile = donor_slide.shapes[4]
        donor_title_style = _first_run_style(donor_tile)
        donor_body_run = donor_tile.text_frame.paragraphs[2].runs[0]
        donor_body_style = (
            str(getattr(donor_body_run.font, "bold", None)),
            str(getattr(donor_body_run.font, "size", None)),
            str(getattr(donor_body_run.font.color, "type", None)),
            str(getattr(donor_body_run.font.color, "theme_color", None)),
        )

        prs = _render_wd(
            {
                "deck_title": "T",
                "slides": [
                    {
                        "type": "wd_mosaic_4",
                        "title": "M4 MAIN",
                        "section_title": "M4 STRAP",
                        "tiles": [
                            {"title": "M4 TILE 1 CAP", "body": "M4 TILE 1 BODY"},
                            {"title": "M4 TILE 2 CAP", "body": "M4 TILE 2 BODY"},
                            {"title": "M4 TILE 3 CAP", "body": "M4 TILE 3 BODY"},
                            {"title": "M4 TILE 4 CAP", "body": "M4 TILE 4 BODY"},
                        ],
                    }
                ],
            }
        )
        out_slide = prs.slides[0]
        # Map donor content tiles after OLE-shape removal in clone: donor 4,5,7,6 -> output 3,4,6,5
        out_tile = out_slide.shapes[3]
        self.assertEqual(out_tile.text_frame.paragraphs[0].text, "M4 TILE 1 CAP")
        self.assertEqual(out_tile.text_frame.paragraphs[2].text, "M4 TILE 1 BODY")
        self.assertEqual(_first_run_style(out_tile), donor_title_style)
        out_body_run = out_tile.text_frame.paragraphs[2].runs[0]
        out_body_style = (
            str(getattr(out_body_run.font, "bold", None)),
            str(getattr(out_body_run.font, "size", None)),
            str(getattr(out_body_run.font.color, "type", None)),
            str(getattr(out_body_run.font.color, "theme_color", None)),
        )
        self.assertEqual(out_body_style, donor_body_style)
        # Numeral shape stays unchanged.
        numeral_texts = [
            (sh.text_frame.text or "").strip()
            for sh in out_slide.shapes
            if hasattr(sh, "text_frame") and (sh.text_frame.text or "").strip().isdigit()
        ]
        self.assertIn("1", numeral_texts)
        self.assertIn("2", numeral_texts)
        self.assertIn("3", numeral_texts)
        self.assertIn("4", numeral_texts)


if __name__ == "__main__":
    unittest.main()

